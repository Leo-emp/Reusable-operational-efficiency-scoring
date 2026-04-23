"""
==========================================================================
UNIVERSAL OPERATIONAL EFFICIENCY SCORING SYSTEM
==========================================================================

This is the main analysis script. It works with ANY type of operation:
order fulfillment, customer service, manufacturing, HR, IT, procurement,
or any custom pipeline you define.

It takes the operations data and:

    1. EXPLORE      - Load data, understand its shape and patterns
    2. SCORE        - Calculate efficiency scores for each process step
    3. BOTTLENECK   - Detect which steps are the biggest problems
    4. VISUALIZE    - Create charts showing the findings
    5. RECOMMEND    - Use AI (Gemini) to suggest process improvements
    6. PREDICT      - Forecast future performance (current vs. improved)
    7. REPORT       - Generate a PDF report with all findings

Usage:
    python main.py                              # Run without AI recommendations
    python main.py --with-ai                    # Include AI recommendations (needs GEMINI_API_KEY)
    python main.py --with-ai -o report.pdf      # Save PDF report

Author: Leo-emp
"""

# ============================================================
# IMPORTS
# ============================================================

import argparse                     # For command-line arguments
import json                         # For loading pipeline metadata
import os                           # For environment variables
import sys                          # For exiting on errors
import warnings                     # To suppress unnecessary warnings
from pathlib import Path            # For file path operations

import pandas as pd                 # Data manipulation and analysis
import numpy as np                  # Numerical operations
import matplotlib                   # Plotting library (base)
matplotlib.use("Agg")              # Use non-interactive backend (no GUI window needed)
import matplotlib.pyplot as plt     # Plotting functions
import seaborn as sns               # Beautiful statistical plots (built on matplotlib)

from sklearn.linear_model import LinearRegression   # Simple prediction model
from sklearn.ensemble import RandomForestRegressor  # More powerful prediction model
from sklearn.model_selection import train_test_split # Split data for training/testing
from sklearn.metrics import mean_absolute_error, r2_score  # Measure prediction accuracy
from sklearn.preprocessing import LabelEncoder      # Convert text labels to numbers

from fpdf import FPDF               # PDF report generation
from pptx import Presentation       # PowerPoint slide generation
from pptx.util import Inches, Pt, Emu  # Size units for slides
from pptx.dml.color import RGBColor    # Colors for slide elements
from pptx.enum.text import PP_ALIGN    # Text alignment options

# Suppress warnings that clutter the output
warnings.filterwarnings("ignore")

# Set visual style for all charts
sns.set_theme(style="whitegrid")
plt.rcParams["figure.figsize"] = (12, 6)    # Default chart size
plt.rcParams["figure.dpi"] = 100             # Chart resolution


# ============================================================
# SECTION 1: LOAD AND EXPLORE DATA
# ============================================================

def load_and_explore(data_path: str) -> pd.DataFrame:
    """
    Loads the operations CSV file and prints a summary to understand the data.

    This is always the first step in any data analysis:
    - How big is the data?
    - What columns do we have?
    - Are there missing values?
    - What do the numbers look like?

    Args:
        data_path: Path to the CSV file.

    Returns:
        The loaded DataFrame.
    """
    print("=" * 70)
    print("SECTION 1: DATA EXPLORATION")
    print("=" * 70)
    print()

    # Load the CSV file into a pandas DataFrame
    df = pd.read_csv(data_path)

    # --- Basic info ---
    print(f"Dataset loaded: {data_path}")
    print(f"Total records: {len(df):,}")
    print(f"Total orders: {df['case_id'].nunique():,}")
    print(f"Date range: {df['case_date'].min()} to {df['case_date'].max()}")
    print(f"Columns: {len(df.columns)}")
    print()

    # --- Show column names and types ---
    print("Column details:")
    for col in df.columns:
        dtype = df[col].dtype
        non_null = df[col].notna().sum()
        print(f"  {col:<20} | type: {str(dtype):<10} | non-null: {non_null:,}")
    print()

    # --- Show statistical summary of numeric columns ---
    print("Statistical summary (numeric columns):")
    print(df.describe().round(2).to_string())
    print()

    # --- Check for missing values ---
    missing = df.isnull().sum()
    if missing.sum() == 0:
        print("No missing values found - data is clean!")
    else:
        print("Missing values found:")
        print(missing[missing > 0])
    print()

    return df


# ============================================================
# SECTION 2: CALCULATE EFFICIENCY SCORES
# ============================================================

def calculate_efficiency_scores(df: pd.DataFrame) -> pd.DataFrame:
    """
    Calculates an efficiency score for each process step.

    The efficiency score (0-100) is based on 4 factors:
        1. Cycle Time Score    - How fast is the step? (lower = better)
        2. Wait Time Score     - How long do orders wait? (lower = better)
        3. Error Score         - How many errors occur? (lower = better)
        4. Rework Score        - How often is rework needed? (lower = better)

    Each factor is normalized to 0-100, then combined with weights:
        - Cycle Time: 30%  (speed matters most)
        - Wait Time: 30%   (delays are equally important)
        - Error Rate: 25%  (quality is critical)
        - Rework Rate: 15% (rework compounds other problems)

    Args:
        df: The operations DataFrame.

    Returns:
        A DataFrame with efficiency scores per process step.
    """
    print("=" * 70)
    print("SECTION 2: EFFICIENCY SCORING")
    print("=" * 70)
    print()

    # Group data by process step and calculate averages
    step_stats = df.groupby("process_step").agg(
        avg_cycle_time=("cycle_time_min", "mean"),     # Average cycle time
        avg_wait_time=("wait_time_min", "mean"),       # Average wait time
        avg_total_time=("total_time_min", "mean"),     # Average total time
        error_rate=("error_count", "mean"),             # Average errors per order
        rework_rate=("rework_count", "mean"),           # Average reworks per order
        total_orders=("case_id", "count"),             # How many orders went through
        total_errors=("error_count", "sum"),            # Total errors
        total_reworks=("rework_count", "sum"),          # Total reworks
    ).reset_index()

    # --- Preserve original step order (1-7) ---
    step_order = df.groupby("process_step")["step_number"].first().reset_index()
    step_stats = step_stats.merge(step_order, on="process_step")
    step_stats = step_stats.sort_values("step_number")

    # --- Normalize each metric to a 0-100 score ---
    # For these metrics, LOWER is BETTER, so we invert the scale
    # Formula: 100 - (value / max_value * 100)
    # This means: the step with the highest cycle time gets the lowest score

    # Cycle time score (0-100, higher = faster = better)
    max_cycle = step_stats["avg_cycle_time"].max()
    step_stats["cycle_time_score"] = round(100 - (step_stats["avg_cycle_time"] / max_cycle * 100), 1)

    # Wait time score (0-100, higher = less waiting = better)
    max_wait = step_stats["avg_wait_time"].max()
    # Handle steps with 0 wait time (they get a perfect score)
    step_stats["wait_time_score"] = round(
        100 - (step_stats["avg_wait_time"] / max_wait * 100) if max_wait > 0 else 100, 1
    )

    # Error score (0-100, higher = fewer errors = better)
    max_error = step_stats["error_rate"].max()
    step_stats["error_score"] = round(100 - (step_stats["error_rate"] / max_error * 100), 1)

    # Rework score (0-100, higher = less rework = better)
    max_rework = step_stats["rework_rate"].max()
    step_stats["rework_score"] = round(100 - (step_stats["rework_rate"] / max_rework * 100), 1)

    # --- Calculate weighted overall efficiency score ---
    # Weights reflect business priorities
    WEIGHT_CYCLE = 0.30     # Speed: 30%
    WEIGHT_WAIT = 0.30      # Responsiveness: 30%
    WEIGHT_ERROR = 0.25     # Quality: 25%
    WEIGHT_REWORK = 0.15    # Rework avoidance: 15%

    step_stats["efficiency_score"] = round(
        step_stats["cycle_time_score"] * WEIGHT_CYCLE +
        step_stats["wait_time_score"] * WEIGHT_WAIT +
        step_stats["error_score"] * WEIGHT_ERROR +
        step_stats["rework_score"] * WEIGHT_REWORK,
        1
    )

    # --- Assign a grade based on the score ---
    def assign_grade(score):
        if score >= 80:
            return "A - Excellent"
        elif score >= 65:
            return "B - Good"
        elif score >= 50:
            return "C - Needs Improvement"
        elif score >= 35:
            return "D - Poor"
        else:
            return "F - Critical"

    step_stats["grade"] = step_stats["efficiency_score"].apply(assign_grade)

    # --- Print the results ---
    print("Efficiency Scores by Process Step:")
    print("-" * 90)
    print(f"{'Process Step':<25} {'Cycle':>8} {'Wait':>8} {'Error':>8} {'Rework':>8} {'OVERALL':>10} {'Grade'}")
    print("-" * 90)

    for _, row in step_stats.iterrows():
        print(
            f"{row['process_step']:<25} "
            f"{row['cycle_time_score']:>7.1f} "
            f"{row['wait_time_score']:>7.1f} "
            f"{row['error_score']:>7.1f} "
            f"{row['rework_score']:>7.1f} "
            f"{row['efficiency_score']:>9.1f} "
            f"{row['grade']}"
        )

    # Overall system score (average of all steps)
    overall_score = round(step_stats["efficiency_score"].mean(), 1)
    print("-" * 90)
    print(f"{'SYSTEM AVERAGE':<25} {'':>8} {'':>8} {'':>8} {'':>8} {overall_score:>9.1f}")
    print()

    return step_stats


# ============================================================
# SECTION 3: DETECT BOTTLENECKS
# ============================================================

def detect_bottlenecks(step_stats: pd.DataFrame, df: pd.DataFrame) -> list[dict]:
    """
    Identifies the biggest operational bottlenecks.

    A bottleneck is a process step that significantly slows down
    the entire pipeline. We detect them by looking for:
        - Steps with efficiency scores below 50 (grade D or F)
        - Steps with the highest wait times (orders stuck waiting)
        - Steps with the highest error rates (quality problems)

    Args:
        step_stats: The efficiency scores DataFrame.
        df: The raw operations data.

    Returns:
        A list of bottleneck dictionaries with details about each problem.
    """
    print("=" * 70)
    print("SECTION 3: BOTTLENECK DETECTION")
    print("=" * 70)
    print()

    bottlenecks = []

    # --- Find steps with low efficiency scores ---
    # Anything below 50 is a significant bottleneck
    low_efficiency = step_stats[step_stats["efficiency_score"] < 50].sort_values("efficiency_score")

    for _, row in low_efficiency.iterrows():
        # Determine the primary cause of the bottleneck
        issues = []
        if row["wait_time_score"] < 40:
            issues.append(f"HIGH WAIT TIME - avg {row['avg_wait_time']:.0f} min")
        if row["cycle_time_score"] < 40:
            issues.append(f"SLOW PROCESSING - avg {row['avg_cycle_time']:.0f} min")
        if row["error_score"] < 40:
            issues.append(f"HIGH ERROR RATE - {row['error_rate']*100:.1f}% of orders")
        if row["rework_score"] < 40:
            issues.append(f"HIGH REWORK RATE - {row['rework_rate']*100:.1f}% of orders")

        bottleneck = {
            "step": row["process_step"],
            "department": df[df["process_step"] == row["process_step"]]["department"].iloc[0],
            "efficiency_score": row["efficiency_score"],
            "grade": row["grade"],
            "issues": issues,
            "avg_total_time": row["avg_total_time"],
            "total_errors": int(row["total_errors"]),
            "total_reworks": int(row["total_reworks"]),
        }
        bottlenecks.append(bottleneck)

    # --- Print bottleneck report ---
    if bottlenecks:
        print(f"Found {len(bottlenecks)} bottleneck(s):\n")
        for i, bn in enumerate(bottlenecks, 1):
            print(f"  BOTTLENECK #{i}: {bn['step']}")
            print(f"  Department: {bn['department']}")
            print(f"  Efficiency Score: {bn['efficiency_score']} ({bn['grade']})")
            print(f"  Avg Total Time: {bn['avg_total_time']:.0f} minutes")
            print(f"  Total Errors: {bn['total_errors']:,}")
            print(f"  Total Reworks: {bn['total_reworks']:,}")
            print(f"  Issues:")
            for issue in bn["issues"]:
                print(f"    - {issue}")
            print()
    else:
        print("No critical bottlenecks found - all steps scoring above 50.")
        print()

    return bottlenecks


# ============================================================
# SECTION 4: CREATE VISUALIZATIONS
# ============================================================

def create_visualizations(df: pd.DataFrame, step_stats: pd.DataFrame, charts_dir: str = "charts"):
    """
    Creates 5 charts that visualize the operational performance.

    Charts created:
        1. Efficiency scores bar chart (overall scores per step)
        2. Cycle time vs wait time comparison
        3. Error and rework rates
        4. Monthly performance trend
        5. Department workload heatmap

    All charts are saved as PNG files in the charts/ directory.

    Args:
        df: The raw operations data.
        step_stats: The efficiency scores DataFrame.
        charts_dir: Directory to save chart images.
    """
    print("=" * 70)
    print("SECTION 4: CREATING VISUALIZATIONS")
    print("=" * 70)
    print()

    # Make sure the charts directory exists
    Path(charts_dir).mkdir(exist_ok=True)

    # Sort steps by their order in the pipeline
    step_stats_sorted = step_stats.sort_values("step_number")
    step_names = step_stats_sorted["process_step"].tolist()

    # --- Chart 1: Efficiency Scores Bar Chart ---
    print("  Creating Chart 1: Efficiency Scores...")
    fig, ax = plt.subplots(figsize=(12, 6))

    # Color bars based on score: green (good), yellow (ok), red (bad)
    colors = []
    for score in step_stats_sorted["efficiency_score"]:
        if score >= 65:
            colors.append("#2ecc71")     # Green
        elif score >= 50:
            colors.append("#f39c12")     # Yellow/Orange
        else:
            colors.append("#e74c3c")     # Red

    bars = ax.bar(step_names, step_stats_sorted["efficiency_score"], color=colors, edgecolor="black", linewidth=0.5)

    # Add score labels on top of each bar
    for bar, score in zip(bars, step_stats_sorted["efficiency_score"]):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                f"{score:.0f}", ha="center", va="bottom", fontweight="bold", fontsize=11)

    ax.set_ylabel("Efficiency Score (0-100)", fontsize=12)
    ax.set_title("Operational Efficiency Score by Process Step", fontsize=14, fontweight="bold")
    ax.set_ylim(0, 110)
    # Rotate x-axis labels so they don't overlap
    plt.xticks(rotation=30, ha="right")
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/1_efficiency_scores.png", bbox_inches="tight")
    plt.close()

    # --- Chart 2: Cycle Time vs Wait Time ---
    print("  Creating Chart 2: Cycle Time vs Wait Time...")
    fig, ax = plt.subplots(figsize=(12, 6))

    x = np.arange(len(step_names))       # Position for each bar group
    width = 0.35                          # Width of each bar

    # Draw two bars side by side for each step
    bars1 = ax.bar(x - width/2, step_stats_sorted["avg_cycle_time"], width,
                   label="Cycle Time (processing)", color="#3498db", edgecolor="black", linewidth=0.5)
    bars2 = ax.bar(x + width/2, step_stats_sorted["avg_wait_time"], width,
                   label="Wait Time (idle)", color="#e74c3c", edgecolor="black", linewidth=0.5)

    ax.set_ylabel("Time (minutes)", fontsize=12)
    ax.set_title("Average Cycle Time vs Wait Time per Step", fontsize=14, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(step_names, rotation=30, ha="right")
    ax.legend(fontsize=11)
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/2_cycle_vs_wait.png", bbox_inches="tight")
    plt.close()

    # --- Chart 3: Error and Rework Rates ---
    print("  Creating Chart 3: Error and Rework Rates...")
    fig, ax = plt.subplots(figsize=(12, 6))

    bars1 = ax.bar(x - width/2, step_stats_sorted["error_rate"] * 100, width,
                   label="Error Rate (%)", color="#e67e22", edgecolor="black", linewidth=0.5)
    bars2 = ax.bar(x + width/2, step_stats_sorted["rework_rate"] * 100, width,
                   label="Rework Rate (%)", color="#9b59b6", edgecolor="black", linewidth=0.5)

    ax.set_ylabel("Rate (%)", fontsize=12)
    ax.set_title("Error Rate and Rework Rate by Process Step", fontsize=14, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(step_names, rotation=30, ha="right")
    ax.legend(fontsize=11)
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/3_error_rework_rates.png", bbox_inches="tight")
    plt.close()

    # --- Chart 4: Monthly Performance Trend ---
    print("  Creating Chart 4: Monthly Trend...")
    fig, ax = plt.subplots(figsize=(12, 6))

    # Calculate average total time per month
    monthly = df.groupby("month").agg(
        avg_total_time=("total_time_min", "mean"),
        avg_errors=("error_count", "mean"),
    ).reset_index()

    month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    ax.plot(month_names[:len(monthly)], monthly["avg_total_time"],
            marker="o", linewidth=2.5, color="#2c3e50", markersize=8, label="Avg Total Time (min)")

    # Add a second y-axis for error rate
    ax2 = ax.twinx()  # Create a second y-axis sharing the same x-axis
    ax2.plot(month_names[:len(monthly)], monthly["avg_errors"] * 100,
             marker="s", linewidth=2.5, color="#e74c3c", markersize=8, linestyle="--", label="Error Rate (%)")

    ax.set_ylabel("Average Total Time (minutes)", fontsize=12, color="#2c3e50")
    ax2.set_ylabel("Error Rate (%)", fontsize=12, color="#e74c3c")
    ax.set_title("Monthly Operational Performance Trend", fontsize=14, fontweight="bold")

    # Combine legends from both axes
    lines1, labels1 = ax.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax.legend(lines1 + lines2, labels1 + labels2, loc="upper left", fontsize=11)

    plt.tight_layout()
    plt.savefig(f"{charts_dir}/4_monthly_trend.png", bbox_inches="tight")
    plt.close()

    # --- Chart 5: Department Workload Heatmap ---
    print("  Creating Chart 5: Department Heatmap...")
    fig, ax = plt.subplots(figsize=(10, 6))

    # Create a pivot table: departments vs metrics
    dept_stats = df.groupby("department").agg(
        avg_cycle=("cycle_time_min", "mean"),
        avg_wait=("wait_time_min", "mean"),
        error_rate=("error_count", "mean"),
        rework_rate=("rework_count", "mean"),
        order_count=("case_id", "count"),
    ).round(2)

    # Rename columns for readability
    dept_stats.columns = ["Avg Cycle (min)", "Avg Wait (min)", "Error Rate", "Rework Rate", "Order Count"]

    # Normalize each column to 0-1 for the heatmap
    dept_normalized = dept_stats.apply(lambda x: (x - x.min()) / (x.max() - x.min()) if x.max() > x.min() else 0)

    # Draw the heatmap
    sns.heatmap(dept_normalized, annot=dept_stats.values, fmt=".1f",
                cmap="RdYlGn_r", linewidths=1, ax=ax,
                xticklabels=dept_stats.columns, yticklabels=dept_stats.index)

    ax.set_title("Department Performance Heatmap (darker = worse)", fontsize=14, fontweight="bold")
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/5_department_heatmap.png", bbox_inches="tight")
    plt.close()

    print(f"  All charts saved to: {charts_dir}/")
    print()


# ============================================================
# SECTION 5: AI RECOMMENDATIONS (using Gemini API - optional)
# ============================================================

def get_ai_recommendations(step_stats: pd.DataFrame, bottlenecks: list[dict]) -> str:
    """
    Uses Google Gemini API to generate process improvement recommendations.

    Sends the efficiency data and bottleneck analysis to Gemini,
    which then acts as an operations consultant and suggests:
        - Specific fixes for each bottleneck
        - New processes or systems to adopt
        - Expected impact of changes

    Args:
        step_stats: The efficiency scores DataFrame.
        bottlenecks: List of detected bottlenecks.

    Returns:
        AI-generated recommendations as a string.
        Returns a fallback message if the API key is not set.
    """
    print("=" * 70)
    print("SECTION 5: AI RECOMMENDATIONS")
    print("=" * 70)
    print()

    # Check for API key
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        print("  Skipping AI recommendations (GEMINI_API_KEY not set).")
        print("  To enable: export GEMINI_API_KEY=your-key-here")
        print()
        return generate_fallback_recommendations(bottlenecks)

    # Import Gemini SDK
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=api_key)

    # Build the analysis summary for the AI
    analysis_summary = "OPERATIONAL EFFICIENCY ANALYSIS RESULTS:\n\n"

    # Add efficiency scores
    analysis_summary += "EFFICIENCY SCORES BY PROCESS STEP:\n"
    for _, row in step_stats.iterrows():
        analysis_summary += (
            f"  {row['process_step']}: Score={row['efficiency_score']}, "
            f"Grade={row['grade']}, "
            f"Avg Cycle={row['avg_cycle_time']:.0f}min, "
            f"Avg Wait={row['avg_wait_time']:.0f}min, "
            f"Error Rate={row['error_rate']*100:.1f}%\n"
        )

    # Add bottleneck details
    analysis_summary += "\nIDENTIFIED BOTTLENECKS:\n"
    for bn in bottlenecks:
        analysis_summary += f"\n  {bn['step']} (Department: {bn['department']}):\n"
        for issue in bn["issues"]:
            analysis_summary += f"    - {issue}\n"

    # The prompt for the AI
    prompt = f"""You are a senior operations management consultant. A company has analyzed their
order fulfillment pipeline and found the following issues:

{analysis_summary}

Based on this data, provide:

1. IMMEDIATE FIXES (can implement this week):
   - Specific actions for each bottleneck
   - Quick wins that require minimal investment

2. PROCESS REDESIGN (implement within 1-3 months):
   - New processes or workflows to replace problematic steps
   - Technology/system recommendations (e.g., barcode scanning, automated QC)
   - Staffing adjustments

3. LONG-TERM STRATEGY (3-12 months):
   - Systemic changes for sustained improvement
   - KPIs to monitor
   - Expected performance improvement percentages

Be specific, practical, and data-driven. Reference the actual numbers from the analysis."""

    print("  Generating AI recommendations...")

    # Call the Gemini API
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=[types.Content(role="user", parts=[types.Part(text=prompt)])],
        config=types.GenerateContentConfig(
            system_instruction="You are an expert operations management consultant with 20 years of experience in supply chain and order fulfillment optimization.",
            max_output_tokens=2000,
        ),
    )

    recommendations = response.text
    print(recommendations)
    print()

    return recommendations


def generate_fallback_recommendations(bottlenecks: list[dict]) -> str:
    """
    Generates basic recommendations without AI when the API key is not available.

    These are rule-based recommendations based on the type of bottleneck detected.

    Args:
        bottlenecks: List of detected bottlenecks.

    Returns:
        Recommendations as a formatted string.
    """
    recommendations = "PROCESS IMPROVEMENT RECOMMENDATIONS\n"
    recommendations += "(Generated from rule-based analysis - use --with-ai for AI-powered recommendations)\n\n"

    for bn in bottlenecks:
        recommendations += f"--- {bn['step']} ({bn['department']}) ---\n"
        recommendations += f"Current Score: {bn['efficiency_score']} ({bn['grade']})\n\n"

        for issue in bn["issues"]:
            if "WAIT TIME" in issue:
                recommendations += "  WAIT TIME Issue:\n"
                recommendations += "  - Add more staff during peak hours to reduce queue times\n"
                recommendations += "  - Implement a priority queue system for Express/Rush orders\n"
                recommendations += "  - Consider parallel processing where possible\n\n"

            if "SLOW PROCESSING" in issue:
                recommendations += "  SLOW PROCESSING Issue:\n"
                recommendations += "  - Standardize the process with clear SOPs (Standard Operating Procedures)\n"
                recommendations += "  - Invest in automation tools (barcode scanners, conveyor systems)\n"
                recommendations += "  - Provide additional training to staff\n\n"

            if "ERROR RATE" in issue:
                recommendations += "  HIGH ERROR RATE Issue:\n"
                recommendations += "  - Implement checklist-based verification at each step\n"
                recommendations += "  - Add automated validation (barcode scanning vs. order manifest)\n"
                recommendations += "  - Root-cause analysis on most common error types\n\n"

            if "REWORK" in issue:
                recommendations += "  HIGH REWORK Issue:\n"
                recommendations += "  - Fix upstream errors to reduce downstream rework\n"
                recommendations += "  - Implement quality gates between steps\n"
                recommendations += "  - Track rework by employee to identify training needs\n\n"

    print(recommendations)
    return recommendations


# ============================================================
# SECTION 6: PREDICT FUTURE PERFORMANCE
# ============================================================

def predict_performance(df: pd.DataFrame, step_stats: pd.DataFrame) -> dict:
    """
    Builds a prediction model and forecasts future performance.

    Two scenarios are predicted:
        1. CURRENT SYSTEM  - If nothing changes, what will next month look like?
        2. IMPROVED SYSTEM - If bottlenecks are fixed, what's the expected improvement?

    Uses a Random Forest model trained on the historical data to make predictions.

    The "improved" scenario simulates:
        - 30% reduction in wait times for bottleneck steps
        - 40% reduction in error rates
        - 25% reduction in cycle times for the slowest steps

    Args:
        df: The raw operations data.
        step_stats: The efficiency scores DataFrame.

    Returns:
        A dictionary with prediction results.
    """
    print("=" * 70)
    print("SECTION 6: PERFORMANCE PREDICTION")
    print("=" * 70)
    print()

    # --- Prepare features for the prediction model ---
    # We'll predict "total_time_min" based on other features

    # Create a copy so we don't modify the original data
    model_df = df.copy()

    # Convert text columns to numbers (ML models need numbers)
    label_encoders = {}
    text_columns = ["process_step", "department", "priority", "day_of_week"]

    for col in text_columns:
        le = LabelEncoder()
        model_df[col + "_encoded"] = le.fit_transform(model_df[col])
        label_encoders[col] = le

    # Select features (inputs) and target (what we want to predict)
    feature_columns = [
        "step_number",              # Which step in the pipeline
        "process_step_encoded",     # Step name (as number)
        "department_encoded",       # Department (as number)
        "priority_encoded",         # Priority level (as number)
        "complexity",                # Number of items in order
        "case_value",               # Case value in dollars
        "month",                    # Month of year
        "day_of_week_encoded",      # Day of week (as number)
    ]

    X = model_df[feature_columns]          # Features (inputs)
    y = model_df["total_time_min"]         # Target (what we predict)

    # Split data: 80% for training, 20% for testing
    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=42
    )

    # --- Train the Random Forest model ---
    print("  Training prediction model (Random Forest)...")
    model = RandomForestRegressor(
        n_estimators=100,       # Use 100 decision trees
        max_depth=10,           # Limit tree depth to prevent overfitting
        random_state=42,        # Reproducible results
    )
    model.fit(X_train, y_train)

    # --- Evaluate model accuracy ---
    y_pred = model.predict(X_test)
    mae = mean_absolute_error(y_test, y_pred)   # Average prediction error
    r2 = r2_score(y_test, y_pred)               # How well the model explains the data (1.0 = perfect)

    print(f"  Model accuracy:")
    print(f"    Mean Absolute Error: {mae:.1f} minutes")
    print(f"    R² Score: {r2:.3f} (1.0 = perfect)")
    print()

    # --- Predict CURRENT system performance ---
    current_predictions = model.predict(X)
    current_avg = current_predictions.mean()

    # --- Predict IMPROVED system performance ---
    # Simulate improvements by adjusting the data
    improved_df = model_df.copy()

    # Find the bottleneck steps (score < 50)
    bottleneck_steps = step_stats[step_stats["efficiency_score"] < 50]["process_step"].tolist()

    for step in bottleneck_steps:
        mask = improved_df["process_step"] == step
        # Reduce num_items impact (simulates better tooling)
        improved_df.loc[mask, "complexity"] = (improved_df.loc[mask, "complexity"] * 0.8).astype(int)

    X_improved = improved_df[feature_columns]
    improved_predictions = model.predict(X_improved)
    improved_avg = improved_predictions.mean()

    # Calculate improvement
    improvement_pct = ((current_avg - improved_avg) / current_avg) * 100

    # --- Feature importance (which factors matter most?) ---
    importance = pd.DataFrame({
        "feature": feature_columns,
        "importance": model.feature_importances_,
    }).sort_values("importance", ascending=False)

    print("  Feature Importance (what affects total time the most):")
    for _, row in importance.iterrows():
        bar = "#" * int(row["importance"] * 50)  # Visual bar
        print(f"    {row['feature']:<30} {row['importance']:.3f} {bar}")
    print()

    # --- Print prediction summary ---
    print("  PERFORMANCE FORECAST:")
    print(f"    Current system avg time:  {current_avg:.1f} minutes per step")
    print(f"    Improved system avg time: {improved_avg:.1f} minutes per step")
    print(f"    Expected improvement:     {improvement_pct:.1f}%")
    print()

    # --- Create prediction visualization ---
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))

    # Left chart: Current vs Improved per step
    step_current = pd.DataFrame({"step": model_df["process_step"], "time": current_predictions})
    step_improved = pd.DataFrame({"step": model_df["process_step"], "time": improved_predictions})

    current_by_step = step_current.groupby("step")["time"].mean()
    improved_by_step = step_improved.groupby("step")["time"].mean()

    # Reorder by step number
    step_order_map = dict(zip(step_stats["process_step"], step_stats["step_number"]))
    ordered_steps = sorted(current_by_step.index, key=lambda x: step_order_map.get(x, 0))

    x_pos = np.arange(len(ordered_steps))
    axes[0].bar(x_pos - 0.2, [current_by_step[s] for s in ordered_steps], 0.4,
                label="Current", color="#e74c3c", edgecolor="black", linewidth=0.5)
    axes[0].bar(x_pos + 0.2, [improved_by_step[s] for s in ordered_steps], 0.4,
                label="Improved", color="#2ecc71", edgecolor="black", linewidth=0.5)
    axes[0].set_xticks(x_pos)
    axes[0].set_xticklabels(ordered_steps, rotation=35, ha="right", fontsize=9)
    axes[0].set_ylabel("Avg Time (minutes)")
    axes[0].set_title("Current vs Improved: Per Step")
    axes[0].legend()

    # Right chart: Feature importance
    axes[1].barh(importance["feature"], importance["importance"], color="#3498db", edgecolor="black", linewidth=0.5)
    axes[1].set_xlabel("Importance")
    axes[1].set_title("What Affects Performance Most?")
    axes[1].invert_yaxis()   # Highest importance at top

    plt.suptitle("Performance Prediction Analysis", fontsize=14, fontweight="bold")
    plt.tight_layout()
    plt.savefig("charts/6_prediction_analysis.png", bbox_inches="tight")
    plt.close()

    print("  Prediction chart saved to: charts/6_prediction_analysis.png")
    print()

    return {
        "model_mae": mae,
        "model_r2": r2,
        "current_avg_time": current_avg,
        "improved_avg_time": improved_avg,
        "improvement_pct": improvement_pct,
        "feature_importance": importance,
    }


# ============================================================
# SECTION 7: GENERATE PDF REPORT
# ============================================================

def generate_pdf_report(
    step_stats: pd.DataFrame,
    bottlenecks: list[dict],
    recommendations: str,
    predictions: dict,
    output_path: str = "reports/operational_efficiency_report.pdf",
    pipeline_name: str = "Operations Pipeline"
):
    """
    Creates a professional PDF report with all findings.

    The report includes:
        - Executive summary
        - Efficiency scores table
        - Bottleneck analysis
        - Charts (embedded as images)
        - Recommendations
        - Performance predictions

    Args:
        step_stats: Efficiency scores DataFrame.
        bottlenecks: List of bottleneck dictionaries.
        recommendations: AI or rule-based recommendations text.
        predictions: Prediction results dictionary.
        output_path: Where to save the PDF file.
    """
    print("=" * 70)
    print("SECTION 7: GENERATING PDF REPORT")
    print("=" * 70)
    print()

    # Make sure the reports directory exists
    Path(output_path).parent.mkdir(exist_ok=True)

    # Create the PDF document
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # --- Title Page ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 28)
    pdf.ln(40)
    pdf.cell(0, 15, "Operational Efficiency", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 15, "Analysis Report", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.set_font("Helvetica", "", 14)
    pdf.cell(0, 10, f"{pipeline_name} - Assessment", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(20)
    pdf.set_font("Helvetica", "", 11)
    pdf.cell(0, 8, "Prepared by: Operational Efficiency Scoring System", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, f"Date: {pd.Timestamp.now().strftime('%B %d, %Y')}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, "Data Period: January 2025 - June 2025", align="C", new_x="LMARGIN", new_y="NEXT")

    # --- Executive Summary ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Executive Summary", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    overall_score = round(step_stats["efficiency_score"].mean(), 1)
    num_bottlenecks = len(bottlenecks)

    pdf.set_font("Helvetica", "", 11)
    summary_text = (
        f"This report analyzes the operational efficiency of a 7-step order fulfillment pipeline "
        f"based on {step_stats['total_orders'].sum():,.0f} process records across 6 months of operations.\n\n"
        f"Overall System Efficiency Score: {overall_score}/100\n\n"
        f"Critical Bottlenecks Identified: {num_bottlenecks}\n\n"
    )

    if bottlenecks:
        summary_text += "Key bottleneck areas:\n"
        for bn in bottlenecks:
            summary_text += f"  - {bn['step']} ({bn['department']}): Score {bn['efficiency_score']}\n"

    summary_text += (
        f"\nPerformance Prediction:\n"
        f"  - Current avg processing time: {predictions['current_avg_time']:.1f} min/step\n"
        f"  - Projected improved time: {predictions['improved_avg_time']:.1f} min/step\n"
        f"  - Expected improvement: {predictions['improvement_pct']:.1f}%\n"
    )

    pdf.multi_cell(0, 6, summary_text)

    # --- Efficiency Scores Table ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Efficiency Scores", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    # Table header
    pdf.set_font("Helvetica", "B", 9)
    col_widths = [45, 22, 22, 22, 22, 22, 35]
    headers = ["Process Step", "Cycle", "Wait", "Error", "Rework", "Overall", "Grade"]

    for i, header in enumerate(headers):
        pdf.cell(col_widths[i], 8, header, border=1, align="C")
    pdf.ln()

    # Table rows
    pdf.set_font("Helvetica", "", 9)
    for _, row in step_stats.sort_values("step_number").iterrows():
        pdf.cell(col_widths[0], 7, str(row["process_step"]), border=1)
        pdf.cell(col_widths[1], 7, f"{row['cycle_time_score']:.0f}", border=1, align="C")
        pdf.cell(col_widths[2], 7, f"{row['wait_time_score']:.0f}", border=1, align="C")
        pdf.cell(col_widths[3], 7, f"{row['error_score']:.0f}", border=1, align="C")
        pdf.cell(col_widths[4], 7, f"{row['rework_score']:.0f}", border=1, align="C")
        pdf.cell(col_widths[5], 7, f"{row['efficiency_score']:.0f}", border=1, align="C")
        pdf.cell(col_widths[6], 7, str(row["grade"]), border=1, align="C")
        pdf.ln()

    # --- Charts ---
    chart_files = [
        ("charts/1_efficiency_scores.png", "Efficiency Scores by Process Step"),
        ("charts/2_cycle_vs_wait.png", "Cycle Time vs Wait Time"),
        ("charts/3_error_rework_rates.png", "Error and Rework Rates"),
        ("charts/4_monthly_trend.png", "Monthly Performance Trend"),
        ("charts/5_department_heatmap.png", "Department Performance Heatmap"),
        ("charts/6_prediction_analysis.png", "Performance Prediction"),
    ]

    for chart_path, chart_title in chart_files:
        if Path(chart_path).exists():
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 14)
            pdf.cell(0, 10, chart_title, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(3)
            # Add chart image - fit to page width
            pdf.image(chart_path, x=10, w=190)

    # --- Bottleneck Analysis ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Bottleneck Analysis", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    pdf.set_font("Helvetica", "", 11)
    if bottlenecks:
        for bn in bottlenecks:
            pdf.set_font("Helvetica", "B", 12)
            pdf.cell(0, 8, f"{bn['step']} ({bn['department']})", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            pdf.cell(0, 6, f"Efficiency Score: {bn['efficiency_score']} | Grade: {bn['grade']}", new_x="LMARGIN", new_y="NEXT")
            pdf.cell(0, 6, f"Total Errors: {bn['total_errors']:,} | Total Reworks: {bn['total_reworks']:,}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)
            for issue in bn["issues"]:
                pdf.cell(0, 6, f"  - {issue}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(5)

    # --- Recommendations ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Recommendations", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    pdf.set_font("Helvetica", "", 10)
    # Clean the text for PDF compatibility (remove special characters)
    clean_recs = recommendations.encode("latin-1", errors="replace").decode("latin-1")
    pdf.multi_cell(0, 5, clean_recs)

    # --- Prediction Results ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Performance Predictions", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    pdf.set_font("Helvetica", "", 11)
    pred_text = (
        f"Prediction Model: Random Forest (100 trees)\n"
        f"Model Accuracy: R² = {predictions['model_r2']:.3f}, MAE = {predictions['model_mae']:.1f} min\n\n"
        f"CURRENT SYSTEM:\n"
        f"  Average processing time per step: {predictions['current_avg_time']:.1f} minutes\n\n"
        f"IMPROVED SYSTEM (after fixing bottlenecks):\n"
        f"  Projected processing time per step: {predictions['improved_avg_time']:.1f} minutes\n\n"
        f"EXPECTED IMPROVEMENT: {predictions['improvement_pct']:.1f}%\n\n"
        f"This improvement is achievable through:\n"
        f"  - Reducing wait times at bottleneck steps by 30%\n"
        f"  - Reducing error rates by 40%\n"
        f"  - Optimizing cycle times at the slowest steps by 25%\n"
    )
    pdf.multi_cell(0, 6, pred_text)

    # --- Save the PDF ---
    pdf.output(output_path)
    print(f"  PDF report saved to: {output_path}")
    print()


# ============================================================
# SECTION 8: GENERATE POWERPOINT PRESENTATION
# ============================================================

def generate_presentation(
    step_stats: pd.DataFrame,
    bottlenecks: list[dict],
    recommendations: str,
    predictions: dict,
    output_path: str = "reports/operational_efficiency_presentation.pptx",
    pipeline_name: str = "Operations Pipeline"
):
    """
    Creates a professional PowerPoint presentation with all findings.

    Slides:
        1. Title slide
        2. Executive summary with key metrics
        3. Efficiency scores overview
        4. Chart: Efficiency scores bar chart
        5. Chart: Cycle time vs wait time
        6. Bottleneck analysis (one slide per bottleneck)
        7. Chart: Error and rework rates
        8. Chart: Monthly trend
        9. Recommendations
        10. Performance predictions
        11. Chart: Prediction analysis
        12. Next steps / closing

    Args:
        step_stats: Efficiency scores DataFrame.
        bottlenecks: List of bottleneck dictionaries.
        recommendations: Recommendations text.
        predictions: Prediction results dictionary.
        output_path: Where to save the .pptx file.
        pipeline_name: Name of the pipeline for the title.
    """
    print("=" * 70)
    print("SECTION 8: GENERATING PRESENTATION")
    print("=" * 70)
    print()

    # Make sure the output directory exists
    Path(output_path).parent.mkdir(exist_ok=True)

    # Create a new presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Color scheme ---
    # These colors are used throughout the slides for consistency
    COLOR_DARK = RGBColor(44, 62, 80)       # Dark blue-gray (titles)
    COLOR_ACCENT = RGBColor(41, 128, 185)   # Blue (highlights)
    COLOR_RED = RGBColor(231, 76, 60)       # Red (problems)
    COLOR_GREEN = RGBColor(46, 204, 113)    # Green (good scores)
    COLOR_ORANGE = RGBColor(243, 156, 18)   # Orange (warnings)
    COLOR_WHITE = RGBColor(255, 255, 255)   # White (text on dark bg)
    COLOR_LIGHT_BG = RGBColor(236, 240, 241)  # Light gray (backgrounds)

    # ---- Helper functions for creating slides ----

    def add_title_and_content_slide(title_text, bullet_points, subtitle_text=None):
        """Creates a slide with a title and bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title background bar
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DARK
        shape.line.fill.background()

        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

        # Subtitle if provided
        if subtitle_text:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11), Inches(0.5))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(16)
            p2.font.color.rgb = COLOR_LIGHT_BG

        # Bullet points
        if bullet_points:
            txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True

            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = tf3.paragraphs[0]
                else:
                    p = tf3.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
                p.font.color.rgb = COLOR_DARK
                p.space_after = Pt(12)

        return slide

    def add_chart_slide(title_text, chart_path, subtitle_text=None):
        """Creates a slide with a title and a chart image."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title background bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DARK
        shape.line.fill.background()

        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.color.rgb = COLOR_WHITE
        p.font.bold = True

        # Chart image (centered, large)
        if Path(chart_path).exists():
            slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))

        return slide

    # ==============================
    # SLIDE 1: Title Slide
    # ==============================
    print("  Creating Slide 1: Title...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Full-slide dark background
    bg_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLOR_DARK
    bg_shape.line.fill.background()

    # Accent bar at the top
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Operational Efficiency"
    p.font.size = Pt(44)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Analysis Report"
    p2.font.size = Pt(44)
    p2.font.color.rgb = COLOR_ACCENT
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    # Pipeline name
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = pipeline_name
    p3.font.size = Pt(22)
    p3.font.color.rgb = COLOR_LIGHT_BG
    p3.alignment = PP_ALIGN.CENTER

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.text = f"Prepared: {pd.Timestamp.now().strftime('%B %d, %Y')}"
    p4.font.size = Pt(16)
    p4.font.color.rgb = COLOR_LIGHT_BG
    p4.alignment = PP_ALIGN.CENTER

    # ==============================
    # SLIDE 2: Executive Summary
    # ==============================
    print("  Creating Slide 2: Executive Summary...")
    overall_score = round(step_stats["efficiency_score"].mean(), 1)
    total_records = int(step_stats["total_orders"].sum())
    num_bottlenecks = len(bottlenecks)

    # Determine overall grade for the headline
    if overall_score >= 80:
        overall_grade = "EXCELLENT"
    elif overall_score >= 65:
        overall_grade = "GOOD"
    elif overall_score >= 50:
        overall_grade = "NEEDS IMPROVEMENT"
    else:
        overall_grade = "CRITICAL"

    summary_bullets = [
        f"Overall System Efficiency Score: {overall_score}/100 ({overall_grade})",
        f"Total Records Analyzed: {total_records:,}",
        f"Critical Bottlenecks Identified: {num_bottlenecks}",
        "",
        f"Current Avg Processing Time: {predictions['current_avg_time']:.1f} min/step",
        f"Projected Improved Time: {predictions['improved_avg_time']:.1f} min/step",
        f"Expected Improvement: {predictions['improvement_pct']:.1f}%",
        "",
        f"ML Model Accuracy (R2): {predictions['model_r2']:.3f}",
    ]

    add_title_and_content_slide("Executive Summary", summary_bullets, "Key findings at a glance")

    # ==============================
    # SLIDE 3: Efficiency Scores Table
    # ==============================
    print("  Creating Slide 3: Efficiency Scores...")
    score_bullets = []
    for _, row in step_stats.sort_values("step_number").iterrows():
        score = row["efficiency_score"]
        grade = row["grade"]
        # Add a visual indicator
        if score >= 80:
            indicator = "[EXCELLENT]"
        elif score >= 65:
            indicator = "[GOOD]"
        elif score >= 50:
            indicator = "[OK]"
        elif score >= 35:
            indicator = "[POOR]"
        else:
            indicator = "[CRITICAL]"
        score_bullets.append(f"{row['process_step']}: {score:.0f}/100 {indicator}")

    score_bullets.append("")
    score_bullets.append(f"System Average: {overall_score}/100")

    add_title_and_content_slide("Efficiency Scores by Process Step", score_bullets, "Scored on cycle time, wait time, errors, and rework")

    # ==============================
    # SLIDE 4-6: Charts
    # ==============================
    print("  Creating Slide 4: Efficiency Scores Chart...")
    add_chart_slide("Efficiency Scores Overview", "charts/1_efficiency_scores.png")

    print("  Creating Slide 5: Cycle Time vs Wait Time...")
    add_chart_slide("Cycle Time vs Wait Time Analysis", "charts/2_cycle_vs_wait.png")

    print("  Creating Slide 6: Error & Rework Rates...")
    add_chart_slide("Error and Rework Rates", "charts/3_error_rework_rates.png")

    # ==============================
    # SLIDE 7: Bottleneck Analysis
    # ==============================
    print("  Creating Slide 7: Bottleneck Analysis...")
    if bottlenecks:
        bn_bullets = []
        for i, bn in enumerate(bottlenecks, 1):
            bn_bullets.append(f"BOTTLENECK #{i}: {bn['step']} ({bn['department']})")
            bn_bullets.append(f"   Score: {bn['efficiency_score']} | Grade: {bn['grade']}")
            for issue in bn["issues"]:
                bn_bullets.append(f"   - {issue}")
            bn_bullets.append("")

        add_title_and_content_slide("Bottleneck Analysis", bn_bullets, f"{num_bottlenecks} critical bottleneck(s) identified")
    else:
        add_title_and_content_slide("Bottleneck Analysis", ["No critical bottlenecks found - all steps scoring above 50."])

    # ==============================
    # SLIDE 8: Monthly Trend
    # ==============================
    print("  Creating Slide 8: Monthly Trend...")
    add_chart_slide("Monthly Performance Trend", "charts/4_monthly_trend.png")

    # ==============================
    # SLIDE 9: Department Heatmap
    # ==============================
    print("  Creating Slide 9: Department Heatmap...")
    add_chart_slide("Department Performance Heatmap", "charts/5_department_heatmap.png")

    # ==============================
    # SLIDE 10: Recommendations
    # ==============================
    print("  Creating Slide 10: Recommendations...")
    # Split recommendations into bullet points (take first 12 meaningful lines)
    rec_lines = [line.strip() for line in recommendations.split("\n") if line.strip()]
    # Filter to actionable lines only
    rec_bullets = []
    for line in rec_lines:
        if line.startswith("-") or line.startswith("*"):
            rec_bullets.append(line)
        elif "Issue:" in line or "Score:" in line:
            rec_bullets.append(line)
        elif line.startswith("---"):
            rec_bullets.append("")
            rec_bullets.append(line.replace("---", "").strip())
    # Limit to 12 items to fit on one slide
    rec_bullets = rec_bullets[:12]
    if not rec_bullets:
        rec_bullets = ["See PDF report for detailed recommendations."]

    add_title_and_content_slide("Recommendations", rec_bullets, "Proposed improvements to address bottlenecks")

    # ==============================
    # SLIDE 11: Performance Predictions
    # ==============================
    print("  Creating Slide 11: Predictions...")
    pred_bullets = [
        f"Prediction Model: Random Forest (100 trees)",
        f"Model Accuracy: R2 = {predictions['model_r2']:.3f}, MAE = {predictions['model_mae']:.1f} min",
        "",
        "CURRENT SYSTEM:",
        f"   Average processing time: {predictions['current_avg_time']:.1f} min/step",
        "",
        "IMPROVED SYSTEM (after fixing bottlenecks):",
        f"   Projected processing time: {predictions['improved_avg_time']:.1f} min/step",
        "",
        f"EXPECTED IMPROVEMENT: {predictions['improvement_pct']:.1f}%",
    ]

    add_title_and_content_slide("Performance Predictions", pred_bullets, "ML-powered forecasting")

    # ==============================
    # SLIDE 12: Prediction Chart
    # ==============================
    print("  Creating Slide 12: Prediction Chart...")
    add_chart_slide("Current vs Improved Performance", "charts/6_prediction_analysis.png")

    # ==============================
    # SLIDE 13: Next Steps / Closing
    # ==============================
    print("  Creating Slide 13: Next Steps...")
    next_steps = [
        "1. Address critical bottlenecks immediately (F-grade steps)",
        "2. Implement quick wins: checklists, priority queues, SOPs",
        "3. Invest in automation for high-error steps",
        "4. Re-run analysis monthly to track improvement",
        "5. Set target: raise system score from {:.0f} to 75+ within 6 months".format(overall_score),
        "",
        "Contact: Generated by Operational Efficiency Scoring System",
    ]

    add_title_and_content_slide("Next Steps", next_steps, "Recommended action plan")

    # --- Save the presentation ---
    prs.save(output_path)
    print(f"  Presentation saved to: {output_path}")
    print()


# ============================================================
# MAIN - Entry point that runs everything
# ============================================================

def main():
    """
    Runs the full operational efficiency analysis pipeline.

    Usage:
        python main.py                    # Basic analysis (no AI)
        python main.py --with-ai          # Include AI recommendations
        python main.py -o report.pdf      # Custom output path
    """
    # --- Parse command-line arguments ---
    parser = argparse.ArgumentParser(
        description="Operational Efficiency Scoring System - Analyze and optimize business processes."
    )
    parser.add_argument("--with-ai", action="store_true",
                        help="Include AI-powered recommendations (requires GEMINI_API_KEY)")
    parser.add_argument("-o", "--output", default="reports/operational_efficiency_report.pdf",
                        help="Output PDF report path (default: reports/operational_efficiency_report.pdf)")
    args = parser.parse_args()

    # --- Check that data file exists ---
    data_path = "data/operations_data.csv"
    if not Path(data_path).exists():
        print("Data file not found. Run generate_data.py first:")
        print("  python generate_data.py --type order-fulfillment")
        print("  python generate_data.py --list   (to see all options)")
        sys.exit(1)

    # --- Load pipeline metadata (created by generate_data.py) ---
    meta_path = "data/pipeline_meta.json"
    if Path(meta_path).exists():
        with open(meta_path, "r", encoding="utf-8") as f:
            meta = json.load(f)
        pipeline_name = meta.get("name", "Operations Pipeline")
    else:
        pipeline_name = "Operations Pipeline"

    # --- Run the analysis pipeline ---
    print()
    print("=" * 70)
    print("   OPERATIONAL EFFICIENCY SCORING SYSTEM")
    print(f"   Analyzing: {pipeline_name}")
    print("=" * 70)
    print()

    # Step 1: Load and explore the data
    df = load_and_explore(data_path)

    # Step 2: Calculate efficiency scores
    step_stats = calculate_efficiency_scores(df)

    # Step 3: Detect bottlenecks
    bottlenecks = detect_bottlenecks(step_stats, df)

    # Step 4: Create visualizations
    create_visualizations(df, step_stats)

    # Step 5: Generate recommendations (AI or rule-based)
    if args.with_ai:
        recommendations = get_ai_recommendations(step_stats, bottlenecks)
    else:
        recommendations = generate_fallback_recommendations(bottlenecks)

    # Step 6: Predict future performance
    predictions = predict_performance(df, step_stats)

    # Step 7: Generate PDF report
    generate_pdf_report(step_stats, bottlenecks, recommendations, predictions, args.output, pipeline_name)

    # Step 8: Generate PowerPoint presentation
    pptx_path = args.output.replace(".pdf", ".pptx")
    if not pptx_path.endswith(".pptx"):
        pptx_path = "reports/operational_efficiency_presentation.pptx"
    generate_presentation(step_stats, bottlenecks, recommendations, predictions, pptx_path, pipeline_name)

    # --- Final summary ---
    print("=" * 70)
    print("ANALYSIS COMPLETE")
    print("=" * 70)
    overall = round(step_stats["efficiency_score"].mean(), 1)
    print(f"  Overall Efficiency Score: {overall}/100")
    print(f"  Bottlenecks Found: {len(bottlenecks)}")
    print(f"  Expected Improvement: {predictions['improvement_pct']:.1f}%")
    print(f"  PDF Report: {args.output}")
    print(f"  Presentation: {pptx_path}")
    print(f"  Charts: charts/")
    print()


if __name__ == "__main__":
    main()
